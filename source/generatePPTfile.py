from pptx import Presentation
from pptx.util import Inches, Pt
from source.readFile import ReadExcelFile

prs = Presentation()

class GeneratePPTFromFile:
    def __init__(self, file_name, excel_file_location):
        self.file_name = file_name
        self.excel_file_location = excel_file_location


    def addSlideToPresentation(self, slide_layout):
        slide_info = prs.slides.add_slide(slide_layout)
        return slide_info
    
    def enterTableData(self, rows, cols, table, sheet_data):
        for i in range(rows):
            for j in range(cols):
                table.cell(i, j).text = str(sheet_data[i][j])

        self.editFontSizeAndFamily(rows, cols, table)
                

    def generatePPTForSingleSheet(self, ppt_location, sheet_name):
        # defined the sheet sizes
        height = Inches(4.5)
        width = Inches(9.52)
        top = Inches(0.1)
        left = Inches(0.1)
        slide_info = self.addSlideToPresentation(1)
        excel_obj = ReadExcelFile(self.excel_file_location)
        rows, cols = excel_obj.rowAndCol(sheet_name)
        # create the table
        table_data = slide_info.shapes.add_table(rows, cols, left, top, width, height).table
        
        # get the sheet data
        sheet_data = excel_obj.readDataFromSheet(sheet_name)
        self.enterTableData(rows, cols, table_data, sheet_data)
        
        # enter slide_data
        prs.save(ppt_location + f'{sheet_name}.pptx')

    def useGenerate(self, sheet_data, rows, cols):
        # defined the sheet sizes
        height = Inches(4.5)
        width = Inches(9.52)
        top = Inches(1)
        left = Inches(1)
        slide_layout = prs.slide_layouts[5]
        slide_info = self.addSlideToPresentation(slide_layout)
        table_data = slide_info.shapes.add_table(rows, cols, left, top, width, height).table
        self.enterTableData(rows, cols, table_data, sheet_data)
        

    def generatePPTForMultipleSheet(self, ppt_location):
        read_object = ReadExcelFile(self.excel_file_location)
        all_sheet_data = read_object.readMultipleSheetData()
        sheet_names = read_object.sheetNames()

        for name, data in all_sheet_data.items():
            rows, cols = read_object.rowAndCol(name)
            self.useGenerate(data, rows, cols)

        prs.save(ppt_location + self.file_name)

    def editFontSizeAndFamily(self, rows, cols, table):
        for i in range(rows):
            for j in range(cols):
                cell = table.cell(i, j)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.name = "EYInterstate Light"

    def __str__(self):
        return self.filename